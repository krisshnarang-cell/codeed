import streamlit as st
import google.generativeai as genai
import os
from gtts import gTTS
import tempfile
import docx
from pptx import Presentation
import speech_recognition as sr
import PyPDF2
from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip
from PIL import Image, ImageDraw, ImageFont
import textwrap
import io
import numpy as np

# =========================
# Configure Gemini API
# =========================
genai.configure(api_key="YOUR_API_KEY")  # 🔑 Replace with your real API key

# =========================
# Custom file size limit (MB)
# =========================
CUSTOM_MAX_SIZE_MB = 1024  # 1 GB

# =========================
# Language Codes
# =========================
lang_codes = {
    "English": "en", "हिंदी": "hi", "Español": "es", "Français": "fr", "Deutsch": "de",
    "中文": "zh", "日本語": "ja", "한국어": "ko", "Русский": "ru", "اردو": "ur",
    "বাংলা": "bn", "తెలుగు": "te", "தமிழ்": "ta", "ગુજરાતી": "gu", "मराठी": "mr",
    "ਪੰਜਾਬੀ": "pa", "ಕನ್ನಡ": "kn", "मारवाड़ी": "rwr", "O‘zbekcha": "uz", "ქართული": "ka",
    "العربية": "ar", "Türkçe": "tr", "ภาษาไทย": "th", "فارسی": "fa", "Shqip": "sq",
    "Nederlands": "nl", "Svenska": "sv", "Italiano": "it", "Việt": "vi", "ລາວ": "lo"
}

# =========================
# Manual translations (for some languages)
# =========================
translations = {
    "en": {"title": "TransformAI","choose_lang": "Choose language:","paste_text": "Paste your text here",
           "upload_file": "Or upload a file (DOCX / PPTX / PDF / Audio)","extra_comments": "Add extra instructions or comments (optional)",
           "output_type": "Choose output type:","generate": "Generate","listen": "🔊 Listen","video": "🎥 Create Video",
           "speech_lang": "Choose speech language:","warning": "Please enter or upload some content.","output": "Output","try_speaker": "Try clicking the speaker below"},
    "hi": {"title": "ट्रांसफॉर्मAI","choose_lang": "भाषा चुनें:","paste_text": "यहाँ अपना पाठ चिपकाएँ",
           "upload_file": "या एक फ़ाइल अपलोड करें (DOCX / PPTX / PDF / ऑडियो)","extra_comments": "अतिरिक्त निर्देश या टिप्पणियाँ जोड़ें (वैकल्पिक)",
           "output_type": "आउटपुट प्रकार चुनें:","generate": "जेनरेट करें","listen": "🔊 सुनें","video": "🎥 वीडियो बनाएँ",
           "speech_lang": "वाक् भाषा चुनें:","warning": "कृपया कुछ सामग्री दर्ज करें या अपलोड करें।","output": "आउटपुट","try_speaker": "नीचे दिए गए स्पीकर पर क्लिक करें"},
    "es": {"title": "TransformAI","choose_lang": "Elige idioma:","paste_text": "Pega tu texto aquí",
           "upload_file": "O sube un archivo (DOCX / PPTX / PDF / Audio)","extra_comments": "Agregue instrucciones o comentarios adicionales (opcional)",
           "output_type": "Elige tipo de salida:","generate": "Generar","listen": "🔊 Escuchar","video": "🎥 Crear Video",
           "speech_lang": "Elige idioma de voz:","warning": "Por favor ingrese o cargue contenido.","output": "Salida","try_speaker": "Haz clic en el altavoz a continuación"},
    "fr": {"title": "TransformAI","choose_lang": "Choisir la langue:","paste_text": "Collez votre texte ici",
           "upload_file": "Ou téléchargez un fichier (DOCX / PPTX / PDF / Audio)","extra_comments": "Ajoutez des instructions ou des commentaires supplémentaires (facultatif)",
           "output_type": "Choisissez le type de sortie:","generate": "Générer","listen": "🔊 Écouter","video": "🎥 Créer Vidéo",
           "speech_lang": "Choisissez la langue de la voix:","warning": "Veuillez entrer ou télécharger du contenu.","output": "Sortie","try_speaker": "Essayez de cliquer sur le haut-parleur ci-dessous"},
    "de": {"title": "TransformAI","choose_lang": "Sprache wählen:","paste_text": "Fügen Sie hier Ihren Text ein",
           "upload_file": "Oder laden Sie eine Datei hoch (DOCX / PPTX / PDF / Audio)","extra_comments": "Fügen Sie zusätzliche Anweisungen oder Kommentare hinzu (optional)",
           "output_type": "Ausgabetyp wählen:","generate": "Generieren","listen": "🔊 Anhören","video": "🎥 Video erstellen",
           "speech_lang": "Wählen Sie die Sprache der Sprachausgabe:","warning": "Bitte geben Sie Inhalte ein oder laden Sie eine Datei hoch.","output": "Ausgabe","try_speaker": "Klicken Sie unten auf den Lautsprecher"}
}

# =========================
# Auto-translate missing UI languages
# =========================
def get_ui_texts(selected_lang):
    if selected_lang in translations:
        return translations[selected_lang]
    english_texts = translations["en"]
    prompt = f"Translate the following UI labels into {selected_lang}:\n\n{english_texts}"
    try:
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        translated_dict = eval(response.candidates[0].content.parts[0].text)
        return translated_dict
    except:
        return english_texts

# =========================
# File Handlers
# =========================
def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    text = []
    for page in reader.pages:
        text.append(page.extract_text() or "")
    return "\n".join(text)

def extract_text_from_audio(file):
    recognizer = sr.Recognizer()
    with sr.AudioFile(file) as source:
        audio_data = recognizer.record(source)
    try:
        return recognizer.recognize_google(audio_data)
    except sr.UnknownValueError:
        return "Could not understand the audio."
    except sr.RequestError:
        return "Speech recognition service error."

# =========================
# Video Helpers
# =========================
def make_slide_image(text, width=1280, height=720, bgcolor=(255,255,255), font_size=36, margin=60):
    img = Image.new("RGB", (width, height), color=bgcolor)
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", font_size)
    except Exception:
        font = ImageFont.load_default()
    wrapped = textwrap.wrap(text, width=40)
    line_height = font.getsize("Ay")[1] + 8
    total_h = line_height * len(wrapped)
    y = max(margin, (height - total_h)//2)
    for line in wrapped:
        w, _ = draw.textsize(line, font=font)
        x = (width - w) // 2
        draw.text((x, y), line, fill=(0,0,0), font=font)
        y += line_height
    return img

def split_text_for_slides(text, max_chars=700):
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    chunks, current = [], ""
    for p in paragraphs:
        if len(current) + len(p) + 1 <= max_chars:
            current = (current + "\n" + p).strip()
        else:
            if current: chunks.append(current)
            if len(p) > max_chars:
                for i in range(0, len(p), max_chars):
                    chunks.append(p[i:i+max_chars])
                current = ""
            else:
                current = p
    if current: chunks.append(current)
    return chunks

# =========================
# UI
# =========================
lang = st.selectbox("Choose language:", list(lang_codes.keys()))
selected_lang = lang_codes[lang]
ui_texts = get_ui_texts(selected_lang)

st.title(ui_texts["title"])
txt = st.text_area(ui_texts["paste_text"])
uploaded_file = st.file_uploader(ui_texts["upload_file"], type=["docx", "pptx", "pdf", "mp3", "wav"])
extra_comments = st.text_area(ui_texts["extra_comments"])
outpt = st.selectbox(ui_texts["output_type"], ["Summary", "Quiz", "Test", "Audio", "Video"])

if "sp_state" not in st.session_state: st.session_state.sp_state = False
if "result" not in st.session_state: st.session_state.result = ""

# =========================
# Process File
# =========================
if uploaded_file:
    uploaded_file.seek(0, os.SEEK_END)
    size_mb = uploaded_file.tell() / (1024*1024)
    uploaded_file.seek(0)
    st.write(f"📂 File size: **{size_mb:.2f} MB**")
    if size_mb > CUSTOM_MAX_SIZE_MB:
        st.error(f"❌ File too large! Limit is {CUSTOM_MAX_SIZE_MB} MB.")
        txt = ""
    else:
        if uploaded_file.type.endswith("document"):
            txt = extract_text_from_docx(uploaded_file)
        elif uploaded_file.type.endswith("presentation"):
            txt = extract_text_from_pptx(uploaded_file)
        elif uploaded_file.type == "application/pdf":
            txt = extract_text_from_pdf(uploaded_file)
        elif uploaded_file.type in ["audio/mpeg", "audio/wav"]:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp:
                tmp.write(uploaded_file.read()); tmp.flush()
                txt = extract_text_from_audio(tmp.name)

# =========================
# Generate
# =========================
if st.button(ui_texts["generate"]):
    if not txt.strip():
        st.warning(ui_texts["warning"])
    else:
        prompt = f"Please generate a {outpt} of the following text in {lang}:\n\n{txt}"
        if extra_comments.strip(): prompt += f"\n\nUser instructions: {extra_comments}"
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        st.session_state.result = response.candidates[0].content.parts[0].text
        st.subheader(ui_texts["output"] + ":")
        st.subheader(ui_texts["try_speaker"])
        st.write(st.session_state.result)

# =========================
# Audio Output
# =========================
if st.session_state.result and st.button(ui_texts["listen"]):
    st.session_state.sp_state = True

if st.session_state.sp_state:
    speech = st.selectbox(ui_texts["speech_lang"], list(lang_codes.keys()))
    if speech != lang:
        translate_prompt = f"Translate the following text to {speech}:\n\n{st.session_state.result}"
        tr_model = genai.GenerativeModel("gemini-1.5-flash")
        tr_response = tr_model.generate_content(translate_prompt)
        text_for_tts = tr_response.candidates[0].content.parts[0].text
    else:
        text_for_tts = st.session_state.result
    tts = gTTS(text=text_for_tts, lang=lang_codes[speech])
    with tempfile.NamedTemporaryFile(delete=True) as tmp:
        tts.save(tmp.name + ".mp3")
        st.audio(tmp.name + ".mp3")

# =========================
# Video Output
# =========================
if st.session_state.result and outpt == "Video":
    st.markdown("### Create Video")
    if st.button(ui_texts["video"]):
        text_to_use = st.session_state.result
        tts_lang = lang_codes.get(lang, "en")
        tts = gTTS(text=text_to_use, lang=tts_lang)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as audio_tmp:
            audio_path = audio_tmp.name; tts.save(audio_path)
        chunks = split_text_for_slides(text_to_use, 700)
        clips = []
        for c in chunks:
            img = make_slide_image(c)
            clips.append(ImageClip(np.array(img)).set_duration(5))
        video = concatenate_videoclips(clips, method="compose")
        video = video.set_audio(AudioFileClip(audio_path))
        output_path = os.path.join(tempfile.gettempdir(), "output.mp4")
        video.write_videofile(output_path, fps=24, codec="libx264", audio_codec="aac", verbose=False, logger=None)
        st.video(output_path)
        with open(output_path, "rb") as f:
            st.download_button("Download Video", f, file_name="output.mp4", mime="video/mp4")


